"""PowerPoint (PPTX) 변환기"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from .base import Asset, ConversionResult, Converter

logger = logging.getLogger(__name__)


class PptxConverter(Converter):
    """PowerPoint PPTX 파일을 Markdown으로 변환하는 변환기"""

    supported_extensions = [".pptx"]

    def convert(self, file_path: Path) -> ConversionResult:
        """
        PPTX 파일을 Markdown으로 변환

        Args:
            file_path: PPTX 파일 경로

        Returns:
            ConversionResult: 변환된 Markdown과 추출된 이미지
        """
        logger.info(f"PPTX 변환 시작: {file_path}")

        markdown_parts: list[str] = []
        assets: list[Asset] = []
        image_counter = 0

        prs = Presentation(str(file_path))
        total_slides = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []

            # 슬라이드 구분선 및 헤더
            slide_title = self._get_slide_title(slide)
            if slide_title:
                slide_parts.append(f"## 슬라이드 {slide_idx}: {slide_title}")
            else:
                slide_parts.append(f"## 슬라이드 {slide_idx}")

            # 슬라이드의 모든 shape 처리
            for shape in slide.shapes:
                # 텍스트 프레임이 있는 shape (제목 제외)
                if shape.has_text_frame:
                    if shape != self._get_title_shape(slide):
                        text = self._extract_text_frame(shape.text_frame)
                        if text.strip():
                            slide_parts.append(text)

                # 테이블 처리
                if shape.has_table:
                    table_md = self._process_table(shape.table)
                    if table_md:
                        slide_parts.append(table_md)

                # 이미지 처리
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_counter += 1
                    asset, img_md = self._process_image(shape, image_counter)
                    if asset:
                        assets.append(asset)
                        slide_parts.append(img_md)

                # 그룹 shape 내부의 이미지 처리
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_shape in shape.shapes:
                        if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_counter += 1
                            asset, img_md = self._process_image(sub_shape, image_counter)
                            if asset:
                                assets.append(asset)
                                slide_parts.append(img_md)

            # 슬라이드 노트 처리
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **노트:** {notes_text}")

            markdown_parts.append("\n\n".join(slide_parts))

            # 슬라이드 간 구분선 (마지막 슬라이드 제외)
            if slide_idx < total_slides:
                markdown_parts.append("\n---\n")

        markdown = "\n".join(markdown_parts)
        markdown = self._cleanup_markdown(markdown)

        logger.info(f"PPTX 변환 완료: {total_slides}개 슬라이드, {len(assets)}개 이미지 추출")

        return ConversionResult(markdown=markdown, assets=assets)

    def _get_title_shape(self, slide):
        """슬라이드의 제목 shape 반환"""
        if slide.shapes.title:
            return slide.shapes.title
        return None

    def _get_slide_title(self, slide) -> str:
        """슬라이드 제목 추출"""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_text_frame(self, text_frame) -> str:
        """텍스트 프레임에서 텍스트 추출 (서식 유지)"""
        paragraphs: list[str] = []

        for para in text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 리스트 레벨 처리 (들여쓰기)
            level = para.level if para.level else 0
            indent = "  " * level

            # 불릿 포인트 처리
            if level > 0 or (para.level is not None and para.level >= 0):
                text = f"{indent}- {text}"

            paragraphs.append(text)

        return "\n".join(paragraphs)

    def _process_table(self, table) -> str:
        """테이블을 Markdown 테이블로 변환"""
        if not table.rows:
            return ""

        rows: list[list[str]] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            rows.append(cells)

        if not rows:
            return ""

        # Markdown 테이블 생성
        md_lines: list[str] = []

        # 헤더 행
        header = rows[0]
        md_lines.append("| " + " | ".join(header) + " |")

        # 구분선
        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        # 데이터 행
        for row in rows[1:]:
            # 열 수가 맞지 않으면 빈 셀로 채움
            while len(row) < len(header):
                row.append("")
            md_lines.append("| " + " | ".join(row[: len(header)]) + " |")

        return "\n".join(md_lines)

    def _process_image(self, shape, counter: int) -> tuple[Asset | None, str]:
        """이미지 shape에서 이미지 추출"""
        try:
            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type

            # 확장자 결정
            ext_map = {
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "image/gif": ".gif",
                "image/bmp": ".bmp",
                "image/tiff": ".tiff",
                "image/webp": ".webp",
            }
            ext = ext_map.get(content_type, ".png")

            filename = f"img_{counter:03d}{ext}"

            asset = Asset(
                filename=filename,
                data=image_bytes,
                mime_type=content_type,
            )

            # 이미지 alt 텍스트 (있으면 사용)
            alt_text = shape.name if shape.name else f"이미지 {counter}"
            img_md = f"![{alt_text}](assets/{filename})"

            return asset, img_md

        except Exception as e:
            logger.warning(f"이미지 추출 실패: {e}")
            return None, ""

    def _cleanup_markdown(self, text: str) -> str:
        """Markdown 텍스트 정리"""
        lines = text.split("\n")
        cleaned_lines: list[str] = []
        empty_count = 0

        for line in lines:
            if not line.strip():
                empty_count += 1
                if empty_count <= 2:
                    cleaned_lines.append(line)
            else:
                empty_count = 0
                cleaned_lines.append(line)

        return "\n".join(cleaned_lines).strip()
